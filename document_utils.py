import io
from typing import Tuple
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

def docx_to_html(docx_data: bytes) -> Tuple[bool, str]:
    """Convert a .docx file to HTML."""
    try:
        doc = DocxDocument(io.BytesIO(docx_data))
        html = ['<div class="docx-preview">']
        
        for para in doc.paragraphs:
            if para.text.strip():
                html.append(f'<p style="margin: 0.5em 0;">{para.text}</p>')
        
        html.append('</div>')
        return True, '\n'.join(html)
    except Exception as e:
        return False, f"Error processing Word document: {str(e)}"

def xlsx_to_html(xlsx_data: bytes) -> Tuple[bool, str]:
    """Convert an Excel file to HTML tables."""
    try:
        wb = load_workbook(filename=io.BytesIO(xlsx_data), read_only=True)
        html = ['<div class="xlsx-preview">']
        
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            html.append(f'<h4>{sheetname}</h4>')
            html.append('<table border="1" cellspacing="0" cellpadding="5" style="border-collapse: collapse; margin-bottom: 20px;">')
            
            for row in ws.iter_rows(values_only=True):
                html.append('<tr>')
                for cell in row:
                    value = str(cell) if cell is not None else ''
                    html.append(f'<td>{value}</td>')
                html.append('</tr>')
            
            html.append('</table>')
        
        html.append('</div>')
        return True, '\n'.join(html)
    except Exception as e:
        return False, f"Error processing Excel file: {str(e)}"

def pptx_to_html(pptx_data: bytes) -> Tuple[bool, str]:
    """Convert a PowerPoint file to HTML."""
    try:
        prs = Presentation(io.BytesIO(pptx_data))
        html = ['<div class="pptx-preview">']
        
        for i, slide in enumerate(prs.slides, 1):
            html.append(f'<div class="slide">')
            html.append(f'<h3>Slide {i}</h3>')
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        html.append(f'<p>{text}</p>')
            
            html.append('</div><hr>')  # Add a separator between slides
        
        html.append('</div>')
        return True, '\n'.join(html)
    except Exception as e:
        return False, f"Error processing PowerPoint file: {str(e)}"

def process_document(file_data: bytes, filename: str) -> Tuple[bool, str]:
    """
    Process a document and return HTML preview.
    Returns: (success: bool, html: str)
    """
    filename = filename.lower()
    
    try:
        if filename.endswith(('.doc', '.docx')):
            return docx_to_html(file_data)
        elif filename.endswith(('.xls', '.xlsx')):
            return xlsx_to_html(file_data)
        elif filename.endswith(('.ppt', '.pptx')):
            return pptx_to_html(file_data)
        else:
            return False, "Unsupported file type"
    except Exception as e:
        return False, f"Error processing document: {str(e)}"
